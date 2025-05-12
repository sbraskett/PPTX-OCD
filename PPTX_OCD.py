from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
import os

def apply_properties_to_target(pptx_path, template_data, font_name):
    prs = Presentation(pptx_path)

    # For each slide in the target presentation
    for slide_idx, slide in enumerate(prs.slides):
        # For each shape in the target slide
        for shape_idx, shape in enumerate(slide.shapes):
            shape_name = shape.name if shape.name else f"Shape_{slide_idx}_{shape_idx}"

            # Check if the shape is in the template data
            if shape_name in template_data:  # Shape name exists in template data
                print(f"Found shape: {shape_name}, Type: {shape.shape_type}")

                # Handle AutoShape (type 1), TextBox (type 17), and Table shapes (type 19)
                # First, change the font name for text-containing shapes
                if shape.shape_type == 1:  # AutoShape
                    print(f"Handling AutoShape (type 1): {shape_name}")
                    # Apply font name to all text in the shape first
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            paragraph.alignment = PP_ALIGN.LEFT
                            for run in paragraph.runs:
                                run.font.name = font_name
                                print(f"Applied font name: {font_name} to run in {shape_name}")

                elif shape.shape_type == 17:  # TextBox (type 17)
                    print(f"Handling TextBox Shape (type 17): {shape_name}")
                    # Apply font name to all text in the TextBox first
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            paragraph.alignment = PP_ALIGN.LEFT
                            for run in paragraph.runs:
                                run.font.name = font_name
                                print(f"Applied font name: {font_name} to run in TextBox")

                elif shape.shape_type == 19:  # Table shape
                    print(f"Handling Table Shape (type 19): {shape_name}")
                    # Apply font name to all text in the table cells first
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if cell.text_frame:  # Check if the cell has a text frame (i.e., it contains text)
                                for paragraph in cell.text_frame.paragraphs:
                                    for run in paragraph.runs:
                                        run.font.name = font_name
                                        print(f"Applied font name: {font_name} to run in table cell")

                # After applying the font, set the position and size
                template_shape = template_data[shape_name]

                # Apply position (top and left) from template
                shape.left = int(template_shape['left'])
                shape.top = int(template_shape['top'])

                # Apply width and height (ensure non-zero sizes)
                if template_shape['width'] == 0:
                    print(f"Width is 0 for {shape_name}, applying default width of 100 points.")
                    shape.width = Pt(100)
                else:
                    shape.width = max(int(template_shape['width']), 1)

                if template_shape['height'] == 0:
                    print(f"Height is 0 for {shape_name}, applying default height of 50 points.")
                    shape.height = Pt(50)
                else:
                    shape.height = max(int(template_shape['height']), 1)

                print(f"Applied Width: {shape.width}, Applied Height: {shape.height}")

    # Save the updated presentation
    prs.save(f"standardized_{os.path.basename(pptx_path)}")


def extract_template_properties(pptx_path):
    prs = Presentation(pptx_path)
    shape_data = {}

    # Iterate over all slides in the presentation
    for slide_idx, slide in enumerate(prs.slides):
        for shape_idx, shape in enumerate(slide.shapes):
            shape_info = {}

            # Skip if the shape does not have a name
            if not shape.name:
                continue

            shape_info['left'] = shape.left
            shape_info['top'] = shape.top
            shape_info['width'] = shape.width
            shape_info['height'] = shape.height

            # Store shape properties by shape name
            shape_data[shape.name] = shape_info

    return shape_data


# Example usage
template_file = "template.pptx"  # Path to the template PowerPoint file
pptx_file = "target.pptx"  # Path to the target PowerPoint file
font_name = "Arial"  # The font name you want to apply to all text

# Step 1: Extract template properties
template_data = extract_template_properties(template_file)

# Step 2: Apply properties to the target presentation
apply_properties_to_target(pptx_file, template_data, font_name)
